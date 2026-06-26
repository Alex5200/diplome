#!/usr/bin/env python3
"""
Генератор презентации — чёрно-белый минималистичный стиль.
Запуск: python docs/generate_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette (black & white) ──
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
GRAY90 = RGBColor(30, 30, 30)
GRAY70 = RGBColor(80, 80, 80)
GRAY50 = RGBColor(130, 130, 130)
GRAY20 = RGBColor(210, 210, 210)
GRAY10 = RGBColor(240, 240, 240)

OUT = "/Users/alexandr/Desktop/vl/new/diplome/docs/diagrams/presentation_bw.pptx"
DIAGRAMS = "/Users/alexandr/Desktop/vl/new/diplome/docs/diagrams"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def _bg(slide, color=BLACK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.line.fill.background()
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    return shp


def _text_box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _para(tf, text, size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT, spacing=1.2):
    p = tf.paragraphs[0] if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "" else tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(size * 0.3)
    p.line_spacing = Pt(size * spacing)
    return p


def _bullet(tf, text, size=16, color=BLACK, level=0, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.level = level
    p.space_after = Pt(4)
    p.line_spacing = Pt(size * 1.3)
    return p


def _img(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        return slide.shapes.add_picture(path, left, top, width, height)
    return None


def _title_bar(slide, text, subtitle=""):
    """Top black bar with title"""
    bar = _rect(slide, 0, 0, W, Inches(1.2), fill_color=BLACK)
    tb = _text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.7))
    _para(tb.text_frame, text, size=32, bold=True, color=WHITE)
    if subtitle:
        tb2 = _text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4))
        _para(tb2.text_frame, subtitle, size=14, color=GRAY20)


def _footer(slide, num):
    bar = _rect(slide, 0, H - Inches(0.4), W, Inches(0.4), fill_color=BLACK)
    tb = _text_box(slide, Inches(0.8), H - Inches(0.35), Inches(3), Inches(0.3))
    _para(tb.text_frame, "Система управления 6-осевым роботом-манипулятором", size=9, color=GRAY50)
    tb2 = _text_box(slide, W - Inches(1.5), H - Inches(0.35), Inches(1), Inches(0.3))
    _para(tb2.text_frame, str(num), size=9, color=GRAY50, align=PP_ALIGN.RIGHT)


def _caption_box(slide, text, left, top, width=Inches(5), height=Inches(0.5), size=11, color=GRAY70):
    tb = _text_box(slide, left, top, width, height)
    _para(tb.text_frame, text, size=size, color=color, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl, BLACK)

# Large centered title
tb = _text_box(sl, Inches(1.5), Inches(1.5), Inches(10), Inches(2))
_para(tb.text_frame, "Система управления", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_para(tb.text_frame, "6-осевым роботом-манипулятором", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_para(tb.text_frame, "на базе сервомоторов ST3215", size=28, color=GRAY50, align=PP_ALIGN.CENTER)

# Horizontal line
_rect(sl, Inches(4), Inches(3.8), Inches(5), Pt(2), fill_color=GRAY50)

tb2 = _text_box(sl, Inches(3), Inches(4.2), Inches(7), Inches(1.5))
_para(tb2.text_frame, "Выпускная квалификационная работа", size=18, color=GRAY20, align=PP_ALIGN.CENTER)
_para(tb2.text_frame, "Ляхов Александр Алексеевич", size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_para(tb2.text_frame, "Московский Политехнический Университет • 2026", size=14, color=GRAY50, align=PP_ALIGN.CENTER)

_footer(sl, 1)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl, WHITE)
_title_bar(sl, "Проблематика", "Почему это актуально?")
_footer(sl, 2)

# Three columns
cols_data = [
    ("Когнитивная перегрузка", "Промышленные фреймворки (ROS2) сложны для обучения. Студенты тратят время на инфраструктуру, а не на кинематику.", "Решение: блочное программирование как «педагогические леса»"),
    ("Задержки управления", "DDS в ROS2 вносит 50-200 мс недетерминированных задержек. Критический путь «зрение → мотор» требует 3-15 мс.", "Решение: прямой UART/RS-485 контроллер"),
    ("Sim-to-Real разрыв", "Тестирование на реальном оборудовании опасно. Нет безопасной среды для отладки и RL-обучения.", "Решение: цифровой двойник MuJoCo"),
]

for i, (title, desc, solution) in enumerate(cols_data):
    x = Inches(0.6 + i * 4.2)
    # Card
    _rect(sl, x, Inches(1.8), Inches(3.8), Inches(5.0), fill_color=GRAY10, line_color=BLACK)
    # Number
    tb = _text_box(sl, x + Inches(0.3), Inches(2.0), Inches(3.5), Inches(0.5))
    _para(tb.text_frame, f"0{i + 1}", size=36, bold=True, color=BLACK)
    # Title
    tb = _text_box(sl, x + Inches(0.3), Inches(2.5), Inches(3.5), Inches(0.5))
    _para(tb.text_frame, title, size=18, bold=True, color=BLACK)
    # Description
    tb = _text_box(sl, x + Inches(0.3), Inches(3.2), Inches(3.5), Inches(2.0))
    _para(tb.text_frame, desc, size=13, color=GRAY70)
    # Solution
    _rect(sl, x + Inches(0.2), Inches(5.3), Inches(3.4), Pt(1), fill_color=BLACK)
    tb = _text_box(sl, x + Inches(0.3), Inches(5.5), Inches(3.5), Inches(1.0))
    _para(tb.text_frame, solution, size=12, bold=True, color=BLACK)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture Overview
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl, WHITE)
_title_bar(sl, "Общая архитектура системы", "Multi-layered — Desktop + Web + ROS2 + AI/ML")
_footer(sl, 3)

# Add the PlantUML diagram
img_path = os.path.join(DIAGRAMS, "architecture_bw.png")
if os.path.exists(img_path):
    _img(sl, img_path, Inches(0.5), Inches(1.5), width=Inches(12.3))
    _caption_box(sl, "Рис. 1 — Многослойная архитектура: пользовательские интерфейсы, сервисы, ядро, коммуникации, аппаратура", Inches(1), Inches(6.6), Inches(11), Inches(0.5))
else:
    tb = _text_box(sl, Inches(1), Inches(2), Inches(11), Inches(4))
    _para(tb.text_frame, "Диаграмма архитектуры", size=24, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    _para(tb.text_frame, "(файл architecture_bw.png)", size=16, color=GRAY50, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Desktop App
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl, WHITE)
_title_bar(sl, "Desktop-приложение", "CustomTkinter — 13 вкладок, 3D-визуализация, блочное программирование")

# Left column - features
tb = _text_box(sl, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.5))
tf = tb.text_frame
_para(tf, "Основные возможности:", size=20, bold=True, color=BLACK)
features = [
    "Подключение USB-Serial, сканирование сервомоторов",
    "Построение карты моторов (ID → сустав)",
    "Jog-режим: ручное управление с контролем скорости",
    "3D-визуализация (DH → forward kinematics, Matplotlib)",
    "Слайдеры углов суставов с обратной связью",
    "Обратная кинематика (IK) — DLS метод",
    "Блочное программирование последовательностей",
    "Teach Pendant — запись и воспроизведение",
    "Позиционные регистры (PR1-100)",
    "AI-управление через VLM (Qwen3 / GPT-4o)",
    "Vision Tracker — PID + VLM трекинг объектов",
    "Emergency Stop, температурные лимиты",
]
for f in features:
    _bullet(tf, f"• {f}", size=14, color=GRAY90)

# Right column - screenshot placeholder
_rect(sl, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.3), fill_color=GRAY10, line_color=BLACK)
tb = _text_box(sl, Inches(7.0), Inches(3.5), Inches(5.5), Inches(1.5))
_para(tb.text_frame, "[Скриншот Desktop]", size=18, color=GRAY50, align=PP_ALIGN.CENTER)
_para(tb.text_frame, "Главное окно приложения", size=13, color=GRAY50, align=PP_ALIGN.CENTER)

_footer(sl, 4)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — 3D Visualization
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl, WHITE)
_title_bar(sl, "3D-визуализация кинематики", "Matplotlib + DH-параметры — визуализация в реальном времени")
_footer(sl, 5)

# Left - screenshot
_rect(sl, Inches(0.5), Inches(1.5), Inches(7.5), Inches(5.5), fill_color=GRAY10, line_color=BLACK)
tb = _text_box(sl, Inches(2.0), Inches(3.5), Inches(5.0), Inches(1.0))
_para(tb.text_frame, "[Скриншот 3D View]", size=18, color=GRAY50, align=PP_ALIGN.CENTER)

# Right features
tb = _text_box(sl, Inches(8.5), Inches(1.5), Inches(4.3), Inches(5.5))
tf = tb.text_frame
_para(tf, "Характеристики:", size=18, bold=True, color=BLACK)
items = [
    "DH-параметры: 6 степеней свободы",
    "Прямая кинематика: матрицы 4×4",
    "Обратная кинематика: DLS (Levenberg-Marquardt)",
    "Слайдеры суставов с автопересчётом",
    "Выбор точки на 3D-сцене (клик)",
    "Waypoints + маршруты",
    "Пресеты проверенных точек",
    "Плавная анимация переходов",
    "Рабочая зона (сфера досягаемости)",
    "Режим Live-синхронизации с моторами",
]
for item in items:
    _bullet(tf, f"• {item}", size=12, color=GRAY90)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Web + ROS2
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl, WHITE)
_title_bar(sl, "Web-интерфейс и ROS2-интеграция")

# Web column
_rect(sl, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5), fill_color=GRAY10, line_color=BLACK)
tb = _text_box(sl, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.0))
tf = tb.text_frame
_para(tf, "Web — React 19 + FastAPI", size=18, bold=True, color=BLACK)
web_items = [
    "Фронтенд: React + TypeScript + Three.js",
    "3D-визуализация робота в браузере",
    "Blockly — блочное программирование",
    "Бэкенд: FastAPI (REST + WebSocket)",
    "JWT-аутентификация",
    "Деплой: Docker Compose",
    "Режим Mock для разработки",
]
for item in web_items:
    _bullet(tf, f"• {item}", size=13, color=GRAY90)

# ROS2 column
_rect(sl, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5), fill_color=GRAY10, line_color=BLACK)
tb = _text_box(sl, Inches(7.1), Inches(1.7), Inches(5.5), Inches(5.0))
tf = tb.text_frame
_para(tf, "ROS2 Humble", size=18, bold=True, color=BLACK)
ros_items = [
    "Ноды: robot_node, monitor_node, IK service",
    "Топики: /robot/joint_states, /joint_cmd",
    "ros2_control Hardware Interface",
    "URDF/xacro-модель робота",
    "Gazebo-симуляция",
    "Rosbag2 — запись/воспроизведение",
    "Критический путь: UART (не DDS)",
]
for item in ros_items:
    _bullet(tf, f"• {item}", size=13, color=GRAY90)

_footer(sl, 6)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AI/ML
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl, WHITE)
_title_bar(sl, "AI/ML подсистема", "VLM-управление, Reinforcement Learning, Computer Vision")
_footer(sl, 7)

# Three columns
ai_cols = [
    ("AIProvider (Strategy)", [
        "Единый фасад: Ollama / LM Studio / OpenAI",
        "VLM: Qwen3 VL, GPT-4o",
        "JSON-команды → управление роботом",
        "Системный промпт с DH-параметрами",
    ]),
    ("Reinforcement Learning", [
        "Среда: RobotArmEnv (Gymnasium)",
        "18-dim наблюдение, 7-dim действие",
        "Агенты: DQN (discrete), PPO (continuous)",
        "Curriculum Learning: reach → pick → place",
        "Sim-to-Real: MuJoCo → реальный робот",
    ]),
    ("Computer Vision", [
        "OpenCV — захват и обработка кадров",
        "ArUco-маркеры: детекция + solvePnP",
        "PID-регулятор для трекинга",
        "YOLOv8 — локальный трекинг (60 FPS)",
        "LeRobot: ACT / Diffusion / pi0 модели",
    ]),
]

for i, (title, items) in enumerate(ai_cols):
    x = Inches(0.5 + i * 4.2)
    _rect(sl, x, Inches(1.5), Inches(3.8), Inches(5.5), fill_color=GRAY10, line_color=BLACK)
    tb = _text_box(sl, x + Inches(0.3), Inches(1.7), Inches(3.5), Inches(0.5))
    _para(tb.text_frame, title, size=16, bold=True, color=BLACK)
    _rect(sl, x + Inches(0.3), Inches(2.3), Inches(3.2), Pt(1), fill_color=BLACK)
    tb = _text_box(sl, x + Inches(0.3), Inches(2.5), Inches(3.5), Inches(4.0))
    tf = tb.text_frame
    for item in items:
        _bullet(tf, f"• {item}", size=12, color=GRAY90)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Key Results
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl, WHITE)
_title_bar(sl, "Результаты работы", "Что разработано?")
_footer(sl, 8)

# Metrics
metrics = [
    ("≈ 15 000", "строк кода Python"),
    ("11", "сервисов приложения"),
    ("13", "вкладок Desktop GUI"),
    ("27", "PlantUML-диаграмм"),
    ("6", "DOF — степеней свободы"),
    ("3-15", "мс latency UART"),
    ("59+", "юнит-тестов"),
    ("4", "интерфейса: Desktop/Web/ROS2/CLI"),
]

for i, (val, desc) in enumerate(metrics):
    col = i % 4
    row = i // 4
    x = Inches(0.6 + col * 3.2)
    y = Inches(1.6 + row * 2.8)

    _rect(sl, x, y, Inches(2.8), Inches(2.3), fill_color=GRAY10, line_color=BLACK)
    tb = _text_box(sl, x, y + Inches(0.3), Inches(2.8), Inches(1.0))
    _para(tb.text_frame, val, size=36, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    tb = _text_box(sl, x, y + Inches(1.3), Inches(2.8), Inches(0.6))
    _para(tb.text_frame, desc, size=13, color=GRAY70, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Architecture Patterns
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl, WHITE)
_title_bar(sl, "Архитектурные паттерны", "Какие паттерны использованы?")
_footer(sl, 9)

patterns = [
    ("EventBus", "Pub-Sub шина событий", "Слабая связанность между компонентами. Авто-отписка через WeakMethod."),
    ("Dependency Injection", "Container", "Ленивая инициализация, singleton/transient, авто-разрешение зависимостей."),
    ("Template Method", "BaseService", "Единый жизненный цикл: init → start → stop. События на каждом этапе."),
    ("Strategy", "AIProvider", "Единый фасад: Ollama / LM Studio / OpenAI. Переключение без изменения кода."),
]

for i, (name, cls, desc) in enumerate(patterns):
    y = Inches(1.5 + i * 1.4)
    _rect(sl, Inches(0.5), y, Inches(2.0), Inches(1.0), fill_color=BLACK)
    tb = _text_box(sl, Inches(0.6), y + Inches(0.15), Inches(1.8), Inches(0.7))
    _para(tb.text_frame, name, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tb = _text_box(sl, Inches(2.7), y + Inches(0.05), Inches(3.0), Inches(0.5))
    _para(tb.text_frame, cls, size=14, bold=True, color=BLACK)
    tb = _text_box(sl, Inches(2.7), y + Inches(0.5), Inches(9.5), Inches(0.5))
    _para(tb.text_frame, desc, size=12, color=GRAY70)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Conclusion
# ════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_bg(sl, BLACK)

tb = _text_box(sl, Inches(1.5), Inches(1.0), Inches(10), Inches(2))
_para(tb.text_frame, "Заключение", size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_rect(sl, Inches(5), Inches(2.8), Inches(3), Pt(2), fill_color=GRAY50)

tb = _text_box(sl, Inches(1.5), Inches(3.3), Inches(10), Inches(3.5))
tf = tb.text_frame
conclusions = [
    "Разработана multi-layered система управления 6-осевым роботом-манипулятором",
    "Единая кодовая база: Desktop + Web + ROS2 + AI/ML — общее ядро",
    "Решены проблемы: когнитивная перегрузка (блоки), latency (UART 3-15ms), Sim-to-Real (MuJoCo)",
    "Интегрированы современные AI/ML подходы: VLM, RL (DQN/PPO), Computer Vision",
    "Открытая архитектура — легко расширять и адаптировать под новые задачи",
]
for c in conclusions:
    _bullet(tf, f"  →  {c}", size=18, color=GRAY20)

# Future
_rect(sl, Inches(0.5), Inches(6.0), W - Inches(1), Pt(1), fill_color=GRAY50)
tb = _text_box(sl, Inches(1.5), Inches(6.3), Inches(10), Inches(0.8))
_para(tb.text_frame, "Перспективы: RRTConnect + FCL (планирование пути) • MediaPipe (жестовое управление) • Zero-copy video (shared memory)", size=12, color=GRAY50, align=PP_ALIGN.CENTER)

_footer(sl, 10)

# ════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Presentation saved: {OUT}")
print(f"  {len(prs.slides)} slides")
